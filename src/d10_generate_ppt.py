"""
D10: Generate defense PPT using python-pptx.
Output: reports/defense_ppt.pptx (~17 slides)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

ROOT = Path(__file__).parent.parent
REPORTS_DIR = ROOT / "reports"

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# Color scheme
BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
ACCENT = RGBColor(0xE8, 0x6A, 0x17)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xDC, 0x35, 0x35)


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, x=0.8, y=0.3, w=11.7, h=0.8, color=WHITE, size=32):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = True
    return tf


def add_body(slide, text, x=0.8, y=1.4, w=11.7, h=5.5, color=WHITE, size=18):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf


def add_img(slide, path, x, y, w=5.5, h=3.5):
    if os.path.exists(path):
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def add_table(slide, data, x, y, w, h, font_size=12):
    """data: list of lists (first row = header)"""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = DARK if i > 0 else WHITE
                if i == 0:
                    paragraph.font.bold = True
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
    return table


def slide_section_header(text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, BLUE)
    add_title(slide, text, x=1, y=3, w=11, h=1.2, size=40)
    return slide


# =========================================================================
# Slide 1: Title
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "基于惯性/磁传感器的人体活动模式识别\n系统设计与实现", x=1, y=1.5, w=11.3, h=2, size=36)
add_body(slide, "福建理工大学  智能科学与技术\n孔子瑜（算法与工程）  李宝平（数据采集与固件）\n指导老师：郑积仕\n2026年6月", x=1, y=4, w=11, h=3, color=LIGHT_GRAY, size=18)

# =========================================================================
# Slide 2: Outline
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "汇报提纲")
add_body(slide, """
  1. 项目概述与目标
  2. 硬件架构与传感器标定
  3. 数据采集与预处理
  4. 特征工程（提取 + 选择 + 降维）
  5. 分类器对比与模型选择
  6. 评估体系与统计检验
  7. 进阶探索（时序平滑 / 1D-CNN / 聚类）
  8. 总结与展望
""", size=20)

# =========================================================================
# Slide 3: Project Overview
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "1. 项目概述")
add_body(slide, """
  ◆ 目标：从传感器数据到模式识别 —— 可穿戴 HAR 全链路工程实践
  ◆ 硬件：ESP32-S3 + MPU6050（Acc+Gyro）+ HMC5883L（Mag）→ 9通道 @50Hz
  ◆ 活动：7类（静坐/站立/步行/跑步/上楼/下楼/跌倒）
  ◆ 被试：2人 × 1次采集 = 14 CSV → 566 滑动窗口
  ◆ 全链路：标定 → 预处理 → 窗口化 → 294维特征 → 6种特征选择 → 13种分类器 → LOSO评估
""", size=20)

add_table(slide, [
    ["阶段", "D1-D2", "D3", "D4-D5", "D6-D7", "D8", "D9"],
    ["内容", "项目启动+标定", "采集+预处理", "特征工程", "分类建模", "评估检验", "进阶探索"],
], x=0.8, y=4.2, w=11.7, h=1.5, font_size=14)

# =========================================================================
# Slide 4: Hardware
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "2. 硬件架构")
add_table(slide, [
    ["组件", "型号/参数", "说明"],
    ["MCU", "ESP32-S3 N16R8", "WiFi/BLE, 16MB Flash, 8MB PSRAM"],
    ["IMU", "MPU6050 @ 0x68", "Acc ±8g, Gyro ±2000dps, I²C"],
    ["磁力计", "HMC5883L @ 0x1E", "3-axis, I²C, 寄存器 X/Z/Y 非标准"],
    ["PCB", "KiCad 2层板", "嘉立创打样, 手工焊接 J1-J4/R1-R2/C1-C2"],
    ["连接", "I²C (SDA=IO17, SCL=IO18)", "4.7kΩ上拉, 100nF去耦"],
    ["指示", "NeoPixel RGB @ GPIO48", "白=编号, 橙=倒计时, 绿=采集, 蓝=OK"],
], x=0.8, y=1.4, w=11.7, h=4.5, font_size=14)

# =========================================================================
# Slide 5: Calibration
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "2. 传感器标定（D2核心）")
add_table(slide, [
    ["传感器", "方法", "模型", "关键指标", "改善"],
    ["加速度计", "六位置静态", "a = M(a_raw-b), 12参数LS", "RMSE 11.35→0.42 m/s²", "96.3%"],
    ["磁力计", "椭球拟合", "DLS→SVD, 硬铁+软铁", "norm_std 7.44→1.22 uT", "76.1%"],
    ["陀螺仪", "Allan方差", "静态零偏+ARW+BI", "零偏[-4.29,-1.21,0.67]°/s", "—"],
], x=0.8, y=1.4, w=11.7, h=2.5, font_size=13)

add_body(slide, """
  标定是模式识别的物理基础："垃圾进，垃圾出"。
  未标定数据的系统误差会让所有特征"失真"，识别率天花板被死死压住。
  加速度计：静止时合加速度模长 = 1g 约束 → 最小二乘求解 12 参数。
  磁力计：3D数据是偏心椭球 → DLS拟合 → SVD分解 → 校正为球。
  陀螺仪：Allan方差 log-log 曲线区分角随机游走(ARW)与零偏不稳定性(BI)。
""", x=0.8, y=4.2, w=11.7, h=3, size=16)

# =========================================================================
# Slide 6: Data Collection
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "3. 数据采集与预处理")
add_table(slide, [
    ["活动", "时长", "窗口数(S01/S02)", "关键特征"],
    ["sit", "65s", "46/46", "加速度≈1g, 角速度≈0"],
    ["stand", "65s", "46/46", "同sit, 姿态角微差"],
    ["walk", "65s", "46/46", "步频~2Hz, 周期性"],
    ["run", "65s", "46/46", "步频~3Hz, 高冲击"],
    ["upstairs", "65s", "46/46", "垂直上推, 偏度为正"],
    ["downstairs", "65s", "46/46", "垂直冲击, 偏度为负"],
    ["fall", "15s", "7/7", "单脉冲, 高幅值"],
], x=0.8, y=1.3, w=8, h=3.8, font_size=12)

add_body(slide, """
  预处理流水线：
  raw CSV → 校准(calib_params.json) → 裁剪首尾2s
  → 去重力(HP 0.3Hz) → 低通(LP 20Hz)
  → 窗口化(2.56s,128点,50%重叠) → (566,128,9)

  窗口标签：多数表决（窗口内取众数）
  数据划分：LOSO按subject_id，标准化在fold内
""", x=9.2, y=1.3, w=3.8, h=5.5, size=14)

# =========================================================================
# Slide 7: Features
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "4. 特征工程：294维 → 292维")
add_table(slide, [
    ["类别", "维度", "内容", "核心判别力"],
    ["时域", "126 (14×9轴)", "mean/std/rms/skew/kurtosis/iqr/zcr...", "运动强度与分布"],
    ["频域", "90 (10×9轴)", "主频/频谱质心/谱熵/频带能量比...", "步频与运动规律"],
    ["模长", "72 (24×3组)", "合Acc/Gyro/Mag的时域+频域", "三维合成不变量"],
    ["复合", "12", "轴间相关/加加速度/垂直水平比/航向std", "跨轴交互与姿态"],
], x=0.8, y=1.3, w=11.7, h=2.8, font_size=13)

add_body(slide, """
  ★ 关键判别特征（按RF重要性排序）：
     ① 合加速度IQR — 运动强度核心指标     ② 垂直/水平能量比 — 上楼vs下楼
     ③ 磁航向标准差 — 转向检测            ④ 加速度skew — 上楼(正) vs 下楼(负)

  ★ 特征选择结论：RFE(RF) 30维最优 (CV 91.7%) > SFS 30维 (90.8%) > ANOVA Top-30 (88.5%)
  ★ PCA: 25维→95%方差, 50维→99%。LDA: 仅6维(7类-1)，高维小样本SVD失效(acc 23.3%)
""", x=0.8, y=4.5, w=11.7, h=2.8, size=15)

# =========================================================================
# Slide 8: Classifiers
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "5. 分类器对比（13种模型 LOSO）")
add_table(slide, [
    ["排名", "模型", "Accuracy", "Macro F1", "分析"],
    ["1", "kNN (k=3,distance)", "0.564", "0.527", "非参数+距离加权, 小样本最优"],
    ["2", "LR (L2)", "0.525", "0.490", "线性基线, 稳健"],
    ["3", "MLP (128,64)", "0.507", "0.472", "小样本下不及kNN"],
    ["4", "GaussianNB", "0.488", "0.442", "特征独立假设过强"],
    ["5", "RBF_SVM", "0.472", "0.430", "核方法需更多样本"],
    ["6", "RandomForest", "0.424", "0.386", "小样本方差大"],
    ["7", "GBDT", "0.281", "0.227", "严重过拟合"],
], x=0.8, y=1.3, w=11.7, h=4.2, font_size=12)

add_body(slide, """
  ★ kNN胜出原因：小样本(每fold仅283训练窗口)下基于实例的方法优于需估计大量参数的模型
  ★ 最小风险贝叶斯（跌倒漏检×10）：跌倒recall 42.9%→57.1%，acc微降0.9pp
""", x=0.8, y=5.8, w=11.7, h=1.2, size=15)

# =========================================================================
# Slide 9: Confusion Matrix
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "6. 评估：混淆矩阵分析")
add_img(slide, REPORTS_DIR / "d8_confusion_matrix.png", x=0.3, y=1.3, w=5.8, h=5.5)
add_body(slide, """
  ★ 核心发现：
     • sit recall = 0% — 全部被错判
       Fold1: sit→run/downstairs（跨被试偏移）
       Fold2: sit→stand 46/46（准静态混淆）
     • walk recall = 93.5% — 最好
     • upstairs recall = 83.7%

  ★ 物理原因：
     sit vs stand: 加速度≈1g, 角速度≈0
     upstairs vs downstairs: 步态相似, 仅垂直方向相反
     fall vs run: 冲击脉冲 vs 周期性峰值

  ★ 统计检验：
     McNemar: kNN vs RBF_SVM p<0.0001（显著）
     Bootstrap 95% CI: [51.2%, 59.4%]
""", x=6.5, y=1.3, w=6.3, h=5.8, size=14)

# =========================================================================
# Slide 10: Evaluation Results Summary
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "6. 评估：完整模型排名")
add_img(slide, REPORTS_DIR / "d8_model_ranking.png", x=0.5, y=1.3, w=6, h=4)
add_img(slide, REPORTS_DIR / "d8_roc_curves.png", x=6.8, y=1.3, w=5.8, h=4)
add_body(slide, """
  左：所有模型 Accuracy + Macro F1 对比（LOSO）。kNN 显著优于其他模型。
  右：OvR ROC 曲线。各类AUC反映一对一可分性，walk AUC最高。
""", x=0.8, y=5.5, w=11.7, h=1.5, size=14)

# =========================================================================
# Slide 11: D9 Temporal Smoothing
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "7. 进阶探索①：时序平滑")
add_img(slide, REPORTS_DIR / "d9_temporal_smoothing_recall.png", x=0.3, y=1.2, w=7, h=3.8)
add_body(slide, """
  方法：
  • 滑窗多数表决 (L=3,5,7)：窗口内多数投票
  • HMM Viterbi解码：训练标签→转移矩阵, kNN概率→发射概率, Viterbi→最优路径

  结果：
  • acc 56.4% → 59.5% (+3.1pp, L=7 MV)
  • walk 93.5% → 100%, upstairs 83.7% → 93.5%
  • sit 恒为 0% — 平滑无法修复特征级混淆

  启示：时序平滑消除偶发抖动有效，但不能创造判别信息。
""", x=7.8, y=1.2, w=5.2, h=5.8, size=13)

# =========================================================================
# Slide 12: D9 CNN + Clustering
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "7. 进阶探索②：1D-CNN + 无监督聚类")
add_img(slide, REPORTS_DIR / "d9_cnn_comparison.png", x=0.3, y=1.2, w=6.3, h=3.2)
add_img(slide, REPORTS_DIR / "d9_clustering_metrics.png", x=6.8, y=1.2, w=6, h=2.8)
add_body(slide, """
  1D-CNN（左）：
  • Accuracy仅25% — 566窗口远不足，严重过拟合
  • sit recall=50%！唯一检测到sit的模型 — 原始信号含判别信息但手工特征未捕获
  • CNN与kNN互补：CNN找到sit找不到stand, kNN反过来

  聚类（右）：
  • PCA降维使ARI翻倍(0.17→0.36) — 292维中约90%为冗余
  • KMeans+PCA: ARI=0.364, NMI=0.516 — 中等聚类质量
  • t-SNE显示walk/run/upstairs成簇, sit/stand高度重叠
""", x=0.8, y=4.6, w=11.7, h=2.8, size=13)

# =========================================================================
# Slide 13: D9 HMM matrix
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "7. 进阶探索③：HMM转移矩阵与聚类可视化")
add_img(slide, REPORTS_DIR / "d9_hmm_transition_matrix.png", x=0.3, y=1.3, w=5.5, h=5.2)
add_img(slide, REPORTS_DIR / "d9_clustering_tsne.png", x=6.2, y=1.3, w=6.5, h=5.5)
add_body(slide, """
  左：HMM转移矩阵 — sit↔stand转移概率高(采集顺序相邻), walk/run自转移概率最高(连续窗口稳定)
  右：t-SNE可视化 — 左上真值, 右上KMeans, 左下GMM, 右下GMM+PCA。
      sit(蓝)与stand(橙)在所有方法中都高度重叠，walk/run/upstairs天然可聚类。
""", x=0.8, y=6.5, w=11.7, h=1, size=11)

# =========================================================================
# Slide 14: Key Results Summary
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "8. 核心结果汇总")

add_table(slide, [
    ["类别", "窗口数", "最佳Recall", "达成方法", "主要混淆"],
    ["walk", "92", "1.000 ★", "kNN+MV L=7", "—"],
    ["upstairs", "92", "0.935 ★", "kNN+MV L=7", "downstairs"],
    ["stand", "92", "0.674", "kNN+MV L=5/7", "sit"],
    ["run", "92", "0.565", "kNN+MV L=7", "walk, up"],
    ["fall", "14", "0.500", "kNN+MV L=5/7", "run"],
    ["downstairs", "92", "0.446", "HMM", "upstairs"],
    ["sit", "92", "0.000 ⚠", "—", "stand, run"],
], x=0.8, y=1.5, w=7, h=3.5, font_size=12)

add_body(slide, """
  ★ 正向成果：
    • 全链路工程闭环（硬件焊接→标定→采集→算法→评估）
    • 6种特征选择+13种模型系统对比 + LOSO严谨评估
    • 时序平滑有效消除偶发抖动（acc +3.1pp）
    • 聚类验证数据结构中等可分（ARI=0.364）

  ★ 诚实局限：
    • LOSO仅2 folds(2被试) — 统计意义有限
    • sit完全不可分 — 跨被试特征偏移（核心瓶颈）
    • CNN极端过拟合(566窗口) — 但sit=50%是有价值信号
""", x=8.2, y=1.5, w=4.8, h=5.5, size=13)

# =========================================================================
# Slide 15: Improvement Roadmap
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "8. 改进路线图")
add_table(slide, [
    ["时间", "方向", "措施", "预期效果"],
    ["短期", "数据增强", "被试≥5人, 每类≥3次, 跌倒≥30s", "LOSO≥5 folds, 统计意义↑"],
    ["短期", "特征优化", "姿态角(pitch/roll), 小波时频特征, 气压计", "sit/stand可分性↑"],
    ["中期", "模型升级", "CNN+手工特征融合, 迁移学习(预训练HAR微调)", "小样本下泛化↑"],
    ["长期", "边缘部署", "模型量化→TF-Lite Micro→ESP32实时推理", "功耗-时延-精度三角"],
], x=0.8, y=1.4, w=11.7, h=3, font_size=13)

add_body(slide, """
  ★ 本次实验最大的"负面结果"反而是最有价值的发现：
    • sit recall=0% → 暴露了跨被试场景下传感器姿态差异的致命影响
    • CNN=25%但sit=50% → 证明了原始信号中有可分辨信息但手工特征遗漏了
    • 聚类ARI=0.36 → 量化了仅靠静态特征的分离上限
  这些"失败"比walk=100%教给我们的东西更多。
""", x=0.8, y=4.8, w=11.7, h=2.5, size=15)

# =========================================================================
# Slide 16: Personal Gains
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_title(slide, "个人收获")

add_body(slide, """
  孔子瑜（算法与工程）：
  • 从"会掉包"到"真懂算法"：理解了每个参数的物理含义和统计前提
  • 掌握了特征工程的物理直觉（为什么上楼偏度为正、下楼为负）
  • 领悟了评估严谨性是区分专业与业余的分水岭（LOSO不是可选项而是必须）
  • 最重要的是：坦诚面对"失败"——sit recall=0% 比 walk=100% 教给我更多

  李宝平（数据采集与固件）：
  • 掌握了 I²C 总线调试的完整流程和 ESP32 MicroPython 开发
  • 通过 PCB SDA/SCL 短路排查理解了硬件调试的系统性思维
  • HMC5883L 寄存器顺序(X,Z,Y非标准)的发现是"手册优先"原则的经典案例
""", size=16)

# =========================================================================
# Slide 17: Thank you
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE)
add_title(slide, "感谢聆听", x=1, y=2.5, w=11, h=1.5, size=48)
add_body(slide, "请各位老师批评指正\n\nGitHub: github.com/kkawo/har-project\n\n孔子瑜  李宝平   |   郑积仕老师指导\n福建理工大学  智能科学与技术   |   2026年6月", x=1, y=4.2, w=11, h=3, size=18, color=LIGHT_GRAY)

# =========================================================================
# Save
# =========================================================================
output_path = REPORTS_DIR / "defense_ppt.pptx"
prs.save(str(output_path))
print(f"PPT saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
