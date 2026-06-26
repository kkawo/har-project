"""
D11: Defense PPT — Error Elimination Focus
Output: 10.D11/答辩PPT.pptx (~25 slides)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent
PROJ = ROOT.parent
R = PROJ / "reports"
F = R / "figures"
OUT = ROOT / "答辩PPT.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
DARK  = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0xAA, 0xAA, 0xBB)
BLUE  = RGBColor(0x3B, 0x82, 0xF6)
ORANGE= RGBColor(0xE8, 0x6A, 0x17)  # accent / problem
GREEN = RGBColor(0x22, 0xC5, 0x5E)  # solution
RED   = RGBColor(0xEF, 0x44, 0x44)  # warning
YELLOW= RGBColor(0xFF, 0xC1, 0x07)

PROBLEM_COLOR = ORANGE
SOLVE_COLOR   = GREEN

# ── Helpers ──
def bg(s, c=DARK):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def tb(s, text, x=0.8, y=0.3, w=11.7, h=0.8, color=WHITE, size=32, bold=True, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = box.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align
    return box.text_frame

def body(s, text, x=0.8, y=1.4, w=11.7, h=5.5, color=WHITE, size=18):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(4)

def slide(title, text):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); tb(s, title); body(s, text); return s

def img_slide(title, img, cap="", w=7.5, h=5, x=2.9, y=1.2):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); tb(s, title)
    try:
        s.shapes.add_picture(str(img), Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception as e:
        body(s, f"[图未找到: {img}]\n{e}", y=2, color=RED)
    if cap:
        tb(s, cap, y=y+h+0.1, size=12, color=GRAY, bold=False, align=PP_ALIGN.CENTER)
    return s

def dual_img_slide(title, img1, img2, cap1="", cap2="", y=1.3, h=4.5):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); tb(s, title)
    w = 5.5
    try:
        s.shapes.add_picture(str(img1), Inches(0.8), Inches(y), Inches(w), Inches(h))
        if cap1: tb(s, cap1, x=0.8, y=y+h+0.1, w=w, size=11, color=GRAY, bold=False, align=PP_ALIGN.CENTER)
    except Exception as e:
        body(s, f"[图未找到: {img1}]", x=0.8, y=y+1, color=RED, size=12)
    try:
        s.shapes.add_picture(str(img2), Inches(7.0), Inches(y), Inches(w), Inches(h))
        if cap2: tb(s, cap2, x=7.0, y=y+h+0.1, w=w, size=11, color=GRAY, bold=False, align=PP_ALIGN.CENTER)
    except Exception as e:
        body(s, f"[图未找到: {img2}]", x=7.0, y=y+1, color=RED, size=12)
    return s

def error_slide(title, problem, root_cause, solution, insight=""):
    """专用：一个工程问题 → 分析 → 解决"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s); tb(s, title)

    # Problem box
    box = s.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(5.5), Inches(1.8))  # rounded rect
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x3D, 0x18, 0x0A)
    box.line.color.rgb = ORANGE; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "问题"; p.font.size = Pt(14); p.font.color.rgb = ORANGE; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = problem; p2.font.size = Pt(15); p2.font.color.rgb = WHITE
    p3 = tf.add_paragraph(); p3.text = f"根因: {root_cause}"; p3.font.size = Pt(13); p3.font.color.rgb = GRAY

    # Solution box
    box2 = s.shapes.add_shape(1, Inches(7.0), Inches(1.3), Inches(5.5), Inches(1.8))
    box2.fill.solid(); box2.fill.fore_color.rgb = RGBColor(0x0A, 0x2E, 0x15)
    box2.line.color.rgb = GREEN; box2.line.width = Pt(1.5)
    tf2 = box2.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]; p.text = "解决"; p.font.size = Pt(14); p.font.color.rgb = GREEN; p.font.bold = True
    p2 = tf2.add_paragraph(); p2.text = solution; p2.font.size = Pt(15); p2.font.color.rgb = WHITE

    if insight:
        body(s, insight, y=3.4, size=14, color=YELLOW)
    return s


# ═══════════════════════════════════════════════════════════════════
# SECTION 1 — 项目概述 (Slides 1-4)
# ═══════════════════════════════════════════════════════════════════

s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, "从传感器数据到模式识别", y=1.5, size=44)
body(s, "可穿戴智能感知系统的全链路工程实践", y=2.3, color=ORANGE, size=28)
body(s, "2024级专周实训  ·  模式识别与统计学习\n\n"
     "团队成员：孔子瑜（算法与工程）  李宝平（数据采集与固件）\n指导老师：郑积仕",
     y=3.8, color=GRAY, size=18)

slide("项目目标",
"""▸ 构建可穿戴设备，实时识别 7 种日常人体活动
  ┃ 静坐 / 站立 / 步行 / 跑步 / 上楼 / 下楼 / 跌倒

▸ 全部在 ESP32-S3 端侧完成（不依赖手机/云端）
  ┃ 传感器：MPU6050 6轴 IMU，腰部佩戴，50Hz
  ┃ 模型：逻辑回归 84维特征 → 7类 Softmax
  ┃ 展示：LED 七色 + HTTP Dashboard + 串口""")

slide("工程全链路",
"""┌──────────┐   ┌──────────┐   ┌──────────┐   ┌──────────┐
│ ① PCB设计  │ → │ ② 嘉立创   │ → │ ③ 手工焊接 │ → │ ④ 传感器标定│
│ KiCad    │   │ 打样     │   │ 阻容/LED  │   │ 六位置法   │
└──────────┘   └──────────┘   └──────────┘   └──────────┘
      ↓
┌──────────┐   ┌──────────┐   ┌──────────┐   ┌──────────┐
│ ⑤ 固件开发  │ → │ ⑥ 数据采集 │ → │ ⑦ 特征工程 │ → │ ⑧ 模型+部署 │
│MicroPython│   │21 trial  │   │ 84维时域  │   │ LR→ESP32 │
└──────────┘   └──────────┘   └──────────┘   └──────────┘
      ↓
┌──────────────────────────────┐
│ ⑨ ESP32 实时推理             │
│ 50Hz → 2.56s窗口 → <2ms推理   │
│ LED + HTTP Dashboard 展示    │
└──────────────────────────────┘""")

slide("硬件系统",
"""▸ 主控：ESP32-S3 N16R8（双核 240MHz, 8MB PSRAM, 16MB Flash）
▸ 传感器：GY-521 (MPU6050) — 3轴加速度 ±8g + 3轴陀螺仪 ±2000°/s
▸ I²C 总线：SDA=IO17, SCL=IO18, Shield PCB 焊接 (R1/R2 4.7kΩ 上拉)
▸ NeoPixel RGB LED @ GPIO48 — 实时活动指示
▸ Shield 扩展板：KiCad 原理图 → Gerber 导出 → 嘉立创打样 → 手工焊接""")


# ═══════════════════════════════════════════════════════════════════
# SECTION 2 — 误差消除 (Slides 5-11) ★核心章节★
# ═══════════════════════════════════════════════════════════════════

# 2.1 误差消除总览
slide("误差来源 & 消除策略总览",
"""▸ 三类误差 + 三级消除：

┌─────────────────┬──────────────────┬──────────────────────┐
│ 误差来源          │ 消除方法           │ 效果                  │
├─────────────────┼──────────────────┼──────────────────────┤
│ ① 传感器系统误差   │ 六位置法标定       │ RMSE 11.35→0.42 m/s² │
│ (零偏/灵敏度/轴耦合)│ 最小二乘拟合       │ 改善 96.3%            │
├─────────────────┼──────────────────┼──────────────────────┤
│ ② 训练-推理不一致  │ 数学表达对齐       │ ZCR修复 / median≈mean │
│ (特征计算差异)     │ 训练数据同步替换    │ IQR≈1.349σ           │
├─────────────────┼──────────────────┼──────────────────────┤
│ ③ 传感器配置漂移   │ 显式寄存器配置      │ 量程统一 ±8g/±2000°/s │
│ (量程/采样率不一致) │ 固件统一初始化      │ 采集=推理 完全一致     │
└─────────────────┴──────────────────┴──────────────────────┘""")

# 2.2 标定图文
img_slide("误差消除① — 传感器六位置标定",
          F / "accel_calib_comparison.png",
          "加速度计六位置法标定 — 每个位置采集5秒，最小二乘求解 bias + scale",
          w=7, h=5, x=3.1, y=1.2)

# 2.3 标定原理详解
slide("六位置法标定 — 原理与实现",
"""▸ 原理：加速度计静止时，三轴合力 = g ≈ 9.80665 m/s²
  ┃ 分别以 6 个朝向静止放置传感器：+Z/-Z/+X/-X/+Y/-Y 各朝下一次
  ┃ 每个位置测得的 (ax, ay, az) 均值应分别等于 (±g, 0, 0), (0, ±g, 0), (0, 0, ±g)

▸ 标定模型:  a_cal = (a_raw - bias) × scale
  ┃ 对每条轴独立做最小二乘: expected = measured / scale + bias → 解出 bias & scale

▸ 结果:
  ┃ 标定前 RMSE = 11.35 m/s²（传感器原始误差 > 1g）
  ┃ 标定后 RMSE = 0.42 m/s²（剩余误差 0.04g）
  ┃ 改善 96.3% — 主要是钟差 bias 被消除""")

# 2.4 预处理管道
img_slide("误差消除② — 数据预处理管道",
          F / "preprocess_overview.png",
          "原始数据 → 标定校准 → 去重力(高通0.3Hz) → 低通滤波(20Hz) → 滑窗(128样本,50%重叠)",
          w=9, h=5.2, x=2.1, y=1.1)

# 2.5 训练-推理特征对齐（关键！）
slide("误差消除③ — 训练与推理的数学表达对齐 ★",
"""这是本项目最核心的工程贡献之一。

▸ 问题：ESP32 MicroPython 无法使用 scipy/numpy，训练用 Python 算的特征和固件算的可能不一致

▸ 三类特征不一致及对齐方案：

  ┃ ① 中位数 (median)：Python 用 np.median() 需要排序 → 固件排序太慢
       解决：训练数据用 mean 替换 median，固件直接算均值

  ┃ ② 四分位距 (IQR)：Python 用 scipy.stats.iqr() 需要排序 → 固件无法实现
       解决：训练数据用 1.349 × std 替换 IQR，固件也直接用 1.349σ

  ┃ ③ 过零率 (ZCR)：Python scipy 先减均值再统计过零 → 固件直接对原始值统计
       后果：加速度含重力 ~9.8，信号永远为正 → 固件 ZCR 永远为 0 → 6个特征全废！
       解决：固件加 mu 参数，先减均值再检测过零""")

# 2.6 ZCR 详解
error_slide("ZCR 过零率 Bug — 最隐蔽的特征不一致",
    "6 个轴的 ZCR 特征在固件端全部为 0，模型收到的输入完全偏离训练分布",
    "Python 训练: scipy 先 x-mean(x) 去掉重力分量，再统计符号变化次数\n固件推理: 直接对原始值检测 (x[i]>=0) != (x[i-1]>=0)\n加速度含重力 ~9.8 → 信号从不跨零 → ZCR 恒为 0",
    "固件 ZCR 加上 mu 参数：_zcr(x, mu) → 先减均值再检测过零\n与 scipy.stats 行为完全一致，6 个 ZCR 特征恢复正确",
    "关键教训: 训练和推理的每一个数学操作必须逐行对齐，不能只对齐\"思路\"")

# 2.7 传感器量程不一致
error_slide("传感器量程 Bug — 推理数据偏 4~8 倍",
    "离线采集固件配置 MPU6050 为 ±8g/±2000°/s，实时推理固件使用默认 ±2g/±250°/s\n但换算系数 ACCL_SENS 和 GYRO_SENS 仍按 ±8g/±2000°/s 计算\n→ 加速度值偏小 4 倍，陀螺值偏小 8 倍 → 特征完全偏离训练分布",
    "采集固件写了 ACCEL_CONFIG (0x1C=0x10) 和 GYRO_CONFIG (0x1B=0x18)\n推理固件只写了 PWR_MGMT1 (0x6B) 唤醒命令，依赖上电默认值\n两个固件在不同量程下运行 → 数据不可比",
    "统一两个固件的 init_mpu() 显式写入全部 5 个配置寄存器：\n0x6B(唤醒) + 0x19(采样率50Hz) + 0x1A(DLPF) + 0x1C(±8g) + 0x1B(±2000°/s)\n采集和推理现在看到完全相同的物理量",
    "关键教训: 嵌入式传感器的配置寄存器必须显式设置，不能依赖默认值")

# 2.8 电源/I²C/WiFi 问题
error_slide("充电宝供电卡死 — I²C & WiFi 工程问题",
    "USB 连电脑一切正常，换充电宝供电后：蓝灯→绿灯→卡死不动",
    "① WiFi 禁用时 socket.bind('0.0.0.0', 80) 在无网络接口下无限阻塞\n② 软重启后 MPU6050 I²C 状态机残留，写寄存器不 ACK (ENODEV)\n③ 充电宝无 USB-CDC 通信，时序与电脑供电不同",
    "① ESP32_IP='0.0.0.0' 时跳过 http_init() — socket 操作全部绕过\n② init_mpu() 自建新 SoftI2C 对象 + read WHO_AM_I 验证 + 3次重试\n③ 加阶段 LED 指示灯 (蓝=MPU, 黄=WiFi, 青=HTTP, 白闪=就绪) 精确定位卡点",
    "关键教训: 嵌入式设备必须考虑脱离调试环境的独立运行场景")


# ═══════════════════════════════════════════════════════════════════
# SECTION 3 — 数据与特征 (Slides 12-14)
# ═══════════════════════════════════════════════════════════════════

slide("数据采集方案",
"""▸ 被试：S01 | 7活动 × 3次重复 = 21 trial | 每次30s（跌倒15s）

▸ 离线采集固件：自动序列 + LED 颜色引导
  蓝=静坐 青=站立 绿=步行 黄=跑步 橙=上楼 紫=下楼 红=跌倒
  常亮=录制中 白闪=完成 彩虹=全部结束

▸ 佩戴：腰带位置，绑紧不滑动 | 方向不重要但采集和推理必须一致

▸ 输出：363 个滑窗 (128样本/窗口, 50%重叠) → CSV 格式""")

img_slide("t-SNE 特征空间可视化",
          R / "d5_tsne.png",
          "84维特征 → t-SNE降维到2D → 7类活动自然聚类，类间边界清晰",
          w=6.5, h=5.5, x=3.4, y=1)

dual_img_slide("特征工程 — 特征选择 & 重要性",
               R / "d5_method_comparison.png",
               R / "d5_rf_importance.png",
               "特征选择方法对比 (Filter/Wrapper/Embedded)",
               "Random Forest 特征重要性 Top 30")


# ═══════════════════════════════════════════════════════════════════
# SECTION 4 — 模型与结果 (Slides 15-19)
# ═══════════════════════════════════════════════════════════════════

slide("模型：逻辑回归 (LogisticRegression)",
"""▸ 多分类策略：One-vs-Rest (OVR) → 7 个二分类器 → Softmax

▸ 输入: 84维手工时域特征 → StandardScaler 标准化
▸ 输出: 7维 logit → Softmax → 概率分布 → argmax
▸ 参数: 84×7=588 权重 + 7 偏置 + 84×2=168 标准化参数

▸ 为什么选 LR 不是 CNN/LSTM？
  ┃ ESP32 MicroPython 跑不动深度学习框架
  ┃ 588次乘法 + Softmax < 2ms
  ┃ 手工特征已提取充足判别信息
  ┃ 可解释 / 可导出 / 可手写循环推理""")

dual_img_slide("模型对比 & 分类器评估",
               R / "d7_model_comparison.png",
               R / "d6_baseline_comparison.png",
               "高级模型 LOSO 精度对比",
               "基线分类器综合对比")

img_slide("混淆矩阵 — 7类分类结果",
          R / "d8_confusion_matrix.png",
          "对角线 = 正确分类 | 非对角线 = 误分类 → 分析最易混淆的活动对",
          w=6, h=5.5, x=3.6, y=1)

dual_img_slide("模型评估 — ROC曲线 & 综合排名",
               R / "d8_roc_curves.png",
               R / "d8_model_ranking.png",
               "各类别 ROC 曲线 (AUC)",
               "模型综合排名 (精度+F1+训练时间)")

slide("模型部署：sklearn → MicroPython 端侧推理",
"""▸ 导出流程：训练完成 → 提取 coef_, intercept_, mean_, scale_
  ┃ → 格式化为 Python 列表 → 写入 model_params.py → ESP32 直接 import

▸ 固件推理代码（核心 10 行，不依赖任何库）：
  # 标准化
  for i in range(84):
      x[i] = (feat[i] - SCALER_MEAN[i]) / SCALER_STD[i]
  # 点积
  for c in range(7):
      s = INTERCEPT[c]
      for i in range(84): s += x[i] * COEF[c][i]
      logits[c] = s
  # Softmax + argmax → 输出类别 + 置信度

▸ 时序平滑：5 帧多数投票 → 消除单帧跳变
▸ 性能：单次推理 < 2ms | 内存占用 ~165KB | 推理频率 1.28s/次""")


# ═══════════════════════════════════════════════════════════════════
# SECTION 5 — 结果展示 (Slides 20-22)
# ═══════════════════════════════════════════════════════════════════

slide("实验结果",
"""▸ 同人自用场景 (S01 本人佩戴)：
  ┃ 训练精度：100%（7类精确率/召回率/F1 全部 1.00）
  ┃ 363 窗口全覆盖，无分类错误

▸ 实时推理性能：
  ┃ 50Hz 采样 → 每 1.28s 输出一次活动识别
  ┃ 单次推理 < 2ms（588 次浮点乘法）
  ┃ 端侧内存占用 ~165KB（含 84×7 权重矩阵）

▸ 展示方式：
  ┃ ESP32 LED NeoPixel 七色实时指示
  ┃ 电脑 Dashboard（串口）+ 手机浏览器（HTTP）
  ┃ /api/status JSON 接口（含传感器波形）

▸ 已知局限：单被试 → 跨人泛化待 S02 补充""")

img_slide("PCA 累计方差 & 降维分析",
          R / "d5_pca_cumvar.png",
          "前N个主成分的累计方差贡献率 → 指导特征降维",
          w=7, h=5, x=3.1, y=1.2)

slide("实时演示 — 现场展示",
"""▸ 硬件准备：ESP32-S3 + Shield + GY-521 + 充电宝 + 腰带

▸ 演示流程：
  1. 上电 → 蓝灯(MPU)→黄灯(WiFi)→青灯(HTTP)→白闪3次(就绪)
  2. 佩戴传感器，做任意动作
  3. ESP32 LED 实时变色：蓝(坐) 青(站) 绿(走) 黄(跑) 橙(上楼) 紫(下楼) 红(跌倒)
  4. 电脑 Dashboard 显示实时波形 + 活动 + 置信度
  5. 手机连热点 → 浏览器 http://<ESP32_IP> → 移动端 Dashboard

▸ 备用方案（无WiFi时）：USB 连接电脑 → 串口 Dashboard""")


# ═══════════════════════════════════════════════════════════════════
# SECTION 6 — 答辩准备 (Slides 23-24)
# ═══════════════════════════════════════════════════════════════════

slide("预判问答 — 常见问题",
"""Q: 100%精度是不是过拟合/数据泄露？
A: 否。这是同人自用场景(self-demo)，训练集包含佩戴者本人，端侧模型就该这么做。

Q: 只有一个被试？跨人泛化呢？
A: 端侧模型为佩戴者本人优化是工业常态（Apple Watch先做个人校准）。S02 数据后续补充。

Q: 为什么不用神经网络而用逻辑回归？
A: ESP32 MicroPython 跑不动深度学习。手工特征+LR 在 2.56s 窗口上已提取到足够判别信息。

Q: 为什么不要磁力计（只用6轴）？
A: 室内磁场干扰大（钢筋、电器），实测贡献小。6轴加速度+陀螺仪已能区分七种活动。

Q: 怎么验证模型学到的不是假特征？
A: 看 t-SNE 聚类效果（7类边界清晰）+ 实时演示（不同姿势做同一动作，输出稳定）""")

# ═══════════════════════════════════════════════════════════════════
# SECTION 7 — 总结 (Slide 25)
# ═══════════════════════════════════════════════════════════════════

s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, "总结", y=0.6, size=36)
body(s, """▸ 从零到一跑通了 硬件→算法→部署 全链路工程实践

▸ 三类误差系统化消除：
  ① 传感器系统误差 → 六位置法标定（RMSE 96.3%改善）
  ② 训练-推理不一致 → 数学表达精确对齐（median/iqr/zcr）
  ③ 配置漂移 → 寄存器显式统一（量程/采样率/滤波器）

▸ 核心能力体现：
  ┃ 遇到问题 → 定位根因 → 工程化解决（不是调包跑模型）
  ┃ 全栈嵌入式 AI：PCB + 固件 + 算法 + 部署

▸ 技术栈：KiCad PCB / MicroPython / scikit-learn / ESP32-S3

▸ 未来方向：多被试扩充 / HMM时序平滑 / 1D-CNN端到端对比""", y=1.4, size=17)

tb(s, "谢谢老师！请提问", y=5.8, size=36, color=YELLOW, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"Done → {OUT} ({len(prs.slides)} slides)")
